import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- Color Palette (Modern Tech) ---
PRIMARY_BLUE = RGBColor(0, 102, 204)  # Bold VectorFlow Blue
TEXT_BLACK = RGBColor(33, 37, 41)  # Soft Black
ACCENT_GRAY = RGBColor(248, 249, 250)  # Off-white / Light Gray
SOFT_BORDER = RGBColor(222, 226, 230)  # Dividers


def add_accent_bar(slide):
    """Adds a stylish vertical accent bar on the left side of the slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_BLUE
    bar.line.fill.background()


def style_header(slide, text):
    """Adds a styled header with a bottom divider line."""
    title_shape = slide.shapes.title
    title_shape.text = text

    # Text Styling
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Add a thin divider line under the title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(12), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = SOFT_BORDER
    line.line.fill.background()


def style_body(slide, content_text):
    """Styles the body text for maximum clarity."""
    body_shape = slide.placeholders[1]
    body_shape.text = content_text
    for paragraph in body_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = TEXT_BLACK
        paragraph.space_before = Pt(15)


# --- Initialize Presentation (16:9) ---
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# 1. TITLE SLIDE
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_accent_bar(slide)
title = slide.shapes.title
title.text = "VectorFlow"
title.text_frame.paragraphs[0].font.size = Pt(88)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE

subtitle = slide.placeholders[1]
subtitle.text = "Multi-Agent 2D Asset Engine\nTeam Clockwork Orange"
subtitle.text_frame.paragraphs[0].font.size = Pt(28)
subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_BLACK

# 2. PROBLEM & USER
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_accent_bar(slide)
style_header(slide, "The Problem: The 'Art Bottleneck'")
style_body(slide,
           "• Solo developers spend 6+ hours manually tracing icons.\n"
           "• Mismatched asset packs ruin game visual consistency.\n"
           "• Current AI tools output unscalable pixels or conversational 'chat'.\n"
           "• Technical Gap: Moving from 'Concept' to 'SVG Code' is a manual chore.")

# 3. WHAT IT DOES
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_accent_bar(slide)
style_header(slide, "The Solution: Agentic Vector Generation")
style_body(slide,
           "• Core: Translates natural language into renderable SVG XML.\n"
           "• Innovation: Two-stage agentic handoff prevents 'instruction drift'.\n"
           "• Speed: Assets generated and cleaned in under 5 seconds.\n"
           "• Outcome: Scalable, mathematically precise production-ready code.")

# 4. LIVE DEMO
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_accent_bar(slide)
style_header(slide, "Live Demo: Scripted Refinement")
style_body(slide,
           "1. Input: 'A cyberpunk neon health potion'.\n"
           "2. Logic: Creative Node generates hex codes and geometry specs.\n"
           "3. Execution: Artist Node streams clean XML to the browser.\n"
           "4. Memory: 'Change liquid to blue' -> Update while keeping bottle shape.")

# 5. ARCHITECTURE
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_accent_bar(slide)
style_header(slide, "Architecture: Reliable Orchestration")
style_body(slide,
           "• Orchestrator: LangGraph Peer-Pipeline state management.\n"
           "• Resilience: 3-Tier Fallback (Gemini -> Claude -> Qwen OSS).\n"
           "• Security: Non-root Docker container with health monitoring.\n"
           "• Protocol: Driver-based persona prevents jailbreaks.")

# 6. MEASUREMENTS
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_accent_bar(slide)
style_header(slide, "Measurements: Evidence of Hardening")
style_body(slide,
           "• Golden Set: 9/10 passing score (Verified by LLM Judge).\n"
           "• Load Test: p50 Latency 3.1s | 2.1 req/s sustained throughput.\n"
           "• FinOps: $0.0005 cost per request via Gemini Flash optimization.\n"
           "• Uptime: Fallback chain successfully handled 50% provider error rates.")

# 7. SAFETY & RELIABILITY
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_accent_bar(slide)
style_header(slide, "Safety & Production Standards")
style_body(slide,
           "• Security: Zero secrets in history; Bearer token auth for MCP.\n"
           "• Control: 10 req/min Rate Limiter protects API credits.\n"
           "• Integrity: Human-in-the-loop gate for expensive Pro models.\n"
           "• Privacy: Prompt hashing and total session isolation.")

# 8. LESSONS & NEXT
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_accent_bar(slide)
style_header(slide, "Lessons and Roadmap")
style_body(slide,
           "• Lesson learned: Infrastructure (sanitizers/fallbacks) > Prompting.\n"
           "• Next Step: Unity/Godot Direct Export Plugin via MCP server.\n"
           "• Thank you! Questions?")

# --- Save to Project Root ---
filename = "VectorFlow_Demo_Day.pptx"
save_path = os.path.join(os.getcwd(), filename)
prs.save(save_path)

print(f"✅ Success! Modern Light-Themed Presentation saved to: {save_path}")